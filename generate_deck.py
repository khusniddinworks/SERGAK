import os
import sys
import re

# Auto-install python-pptx if missing
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("Installing python-pptx library...")
    os.system('pip install python-pptx')
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

def parse_slides(md_path):
    with open(md_path, 'r', encoding='utf-8') as f:
        content = f.read()
    
    # Split slides by Markdown horizontal rules
    raw_slides = content.split('---')
    slides_data = []
    
    for rs in raw_slides:
        rs = rs.strip()
        if not rs or rs.startswith('# SERGAK'):
            continue
        
        lines = rs.split('\n')
        slide_info = {
            'page_title': '',
            'sarlavha': '',
            'vazifa': '',
            'natija': '',
            'points': []
        }
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
            
            # Match Markdown headers like "## Sahifa X: ..."
            if line.startswith('##'):
                slide_info['page_title'] = line.replace('##', '').strip()
            # Match "**Sarlavha:** ..."
            elif 'Sarlavha:' in line:
                slide_info['sarlavha'] = re.sub(r'\*\*Sarlavha:\*\*', '', line).strip()
            # Match "**Vazifa:** ..."
            elif 'Vazifa:' in line:
                slide_info['vazifa'] = re.sub(r'\*\*Vazifa:\*\*', '', line).strip()
            # Match "**Natija (Amaliy ish):** ..."
            elif 'Natija' in line:
                slide_info['natija'] = re.sub(r'\*\*Natija \(Amaliy ish\):\*\*', '', line).strip()
                slide_info['natija'] = re.sub(r'\*\*Natija:\*\*', '', slide_info['natija']).strip()
            # Match bullet points
            elif line.startswith('*') or line.startswith('-'):
                clean_point = re.sub(r'^[\*\-\s]+', '', line).strip()
                # Remove bold markers from bullets if present
                clean_point = clean_point.replace('**', '')
                slide_info['points'].append(clean_point)
                
        if slide_info['page_title'] or slide_info['sarlavha']:
            slides_data.append(slide_info)
            
    return slides_data

def create_presentation(slides_data, output_path):
    prs = Presentation()
    
    # Widescreen 16:9 dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Use blank layout
    blank_layout = prs.slide_layouts[6]
    
    # Premium Dark Color Palette
    bg_color = RGBColor(15, 23, 42)      # Deep Navy/Slate (#0F172A)
    accent_green = RGBColor(0, 230, 118)  # Glowing Neon Green (#00E676)
    text_white = RGBColor(248, 250, 252) # Clean Off-White (#F8FAFC)
    text_gray = RGBColor(148, 163, 184)  # Muted Slate Gray (#94A3B8)
    
    for idx, data in enumerate(slides_data):
        slide = prs.slides.add_slide(blank_layout)
        
        # Apply solid background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color
        
        # Slide Index / Indicator at top-right
        indicator_box = slide.shapes.add_textbox(Inches(11.0), Inches(0.4), Inches(2.0), Inches(0.8))
        tf_ind = indicator_box.text_frame
        p_ind = tf_ind.paragraphs[0]
        p_ind.text = f"{idx + 1} / {len(slides_data)}"
        p_ind.font.name = 'Inter'
        p_ind.font.size = Pt(12)
        p_ind.font.bold = True
        p_ind.font.color.rgb = text_gray
        p_ind.alignment = PP_ALIGN.RIGHT
        
        # 1. Main Title Box
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.5), Inches(1.5))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        
        p_page = tf_title.paragraphs[0]
        p_page.text = data['page_title'].upper()
        p_page.font.name = 'Inter'
        p_page.font.size = Pt(14)
        p_page.font.bold = True
        p_page.font.color.rgb = text_gray
        p_page.space_after = Pt(4)
        
        p_main = tf_title.add_paragraph()
        p_main.text = data['sarlavha'] if data['sarlavha'] else "SERGAK AI"
        p_main.font.name = 'Inter'
        p_main.font.size = Pt(28)
        p_main.font.bold = True
        p_main.font.color.rgb = accent_green
        
        # 2. Content Box
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(11.7), Inches(4.5))
        tf_content = content_box.text_frame
        tf_content.word_wrap = True
        
        # Determine content format
        if data['vazifa'] or data['natija']:
            # Project Days Format (Vazifa & Natija)
            if data['vazifa']:
                p_vaz_lbl = tf_content.paragraphs[0]
                p_vaz_lbl.text = "🎯 BUGUNGI VAZIFA:"
                p_vaz_lbl.font.name = 'Inter'
                p_vaz_lbl.font.size = Pt(15)
                p_vaz_lbl.font.bold = True
                p_vaz_lbl.font.color.rgb = accent_green
                p_vaz_lbl.space_after = Pt(2)
                
                p_vaz_val = tf_content.add_paragraph()
                p_vaz_val.text = data['vazifa']
                p_vaz_val.font.name = 'Inter'
                p_vaz_val.font.size = Pt(17)
                p_vaz_val.font.color.rgb = text_white
                p_vaz_val.space_after = Pt(20)
                
            if data['natija']:
                p_nat_lbl = tf_content.add_paragraph()
                p_nat_lbl.text = "🚀 AMALIY ISH VA NATIJA:"
                p_nat_lbl.font.name = 'Inter'
                p_nat_lbl.font.size = Pt(15)
                p_nat_lbl.font.bold = True
                p_nat_lbl.font.color.rgb = accent_green
                p_nat_lbl.space_after = Pt(2)
                
                p_nat_val = tf_content.add_paragraph()
                p_nat_val.text = data['natija']
                p_nat_val.font.name = 'Inter'
                p_nat_val.font.size = Pt(17)
                p_nat_val.font.color.rgb = text_white
                p_nat_val.space_after = Pt(10)
        else:
            # Bullet point format (Intro slides)
            first_bullet = True
            for pt in data['points']:
                p_pt = tf_content.paragraphs[0] if first_bullet else tf_content.add_paragraph()
                first_bullet = False
                p_pt.text = f"•  {pt}"
                p_pt.font.name = 'Inter'
                p_pt.font.size = Pt(18)
                p_pt.font.color.rgb = text_white
                p_pt.space_after = Pt(14)
                
        # Footer
        footer_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.8), Inches(6.0), Inches(0.4))
        tf_foot = footer_box.text_frame
        p_foot = tf_foot.paragraphs[0]
        p_foot.text = "SERGAK - O'zbekistonning birinchi 100% oflayn AI kiberxavfsizlik ilovasi"
        p_foot.font.name = 'Inter'
        p_foot.font.size = Pt(10)
        p_foot.font.color.rgb = text_gray

    prs.save(output_path)
    print(f"Presentation saved successfully at: {output_path}")

if __name__ == '__main__':
    md_file = r"C:\Users\ki770\.gemini\antigravity\brain\3f347020-3a97-4525-8fab-d2f3ec5cbda7\artifacts\presentation_slides.md"
    output_file = r"C:\Users\ki770\OneDrive\Desktop\py_files\SERGAK_PitchDeck.pptx"
    
    if not os.path.exists(md_file):
        print(f"Error: Markdown file not found at {md_file}")
        sys.exit(1)
        
    slides = parse_slides(md_file)
    create_presentation(slides, output_file)
